#!/usr/bin/env python3
"""
ScaleFlow Branded Presentation Generator

Generates visually branded .pptx presentations from JSON input.
Follows the same pattern as shared/generate_branded_docx.py but targets PowerPoint.

Usage:
    python3 generate_pptx.py --input data.json --output output.pptx

Supported slide layouts:
    - title           Big title + subtitle, brand colors, logo, dark background
    - section_header  Section divider, large centered text, accent color bar
    - content         Title + bullet points body
    - image           Title + image (or placeholder text if image missing)
    - comparison      Title + data table (alias: table)
    - two_column      Title + two content columns
    - closing         Thank you slide with brand presence

JSON input structure: see Deck Creator SKILL.md for full schema.
"""

import json
import sys
import os
import argparse

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image


# ============================================================
# COLOR UTILITIES
# ============================================================

def hex_to_rgb(hex_color):
    """Convert a hex color string to a python-pptx RGBColor object."""
    hex_color = hex_color.lstrip("#")
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


def hex_to_tuple(hex_color):
    """Convert a hex color string to an (r, g, b) tuple."""
    hex_color = hex_color.lstrip("#")
    return (
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


def lighten_hex(hex_color, factor=0.92):
    """Lighten a hex color by blending toward white."""
    hex_color = hex_color.lstrip("#")
    r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    r = int(r + (255 - r) * factor)
    g = int(g + (255 - g) * factor)
    b = int(b + (255 - b) * factor)
    return f"#{r:02X}{g:02X}{b:02X}"


def darken_hex(hex_color, factor=0.3):
    """Darken a hex color by reducing brightness."""
    hex_color = hex_color.lstrip("#")
    r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    r = int(r * (1 - factor))
    g = int(g * (1 - factor))
    b = int(b * (1 - factor))
    return f"#{r:02X}{g:02X}{b:02X}"


# ============================================================
# FONT MAPPING — always use system-safe fallbacks
# ============================================================

FONT_FALLBACKS = {
    "Bebas Neue": "Arial Black",
    "Inter": "Calibri",
    "Montserrat": "Arial",
    "Playfair Display": "Georgia",
    "Roboto": "Calibri",
    "DM Sans": "Calibri",
    "Poppins": "Calibri",
    "Helvetica Neue": "Helvetica",
    "Garamond": "Garamond",
}

STYLE_FONT_MAP = {
    "modern": ("Arial", "Calibri"),
    "classic": ("Georgia", "Georgia"),
    "bold": ("Arial Black", "Arial"),
    "minimal": ("Helvetica", "Calibri"),
    "editorial": ("Garamond", "Garamond"),
    "glassmorphism": ("Arial Black", "Calibri"),
}


# ============================================================
# SLIDE DIMENSIONS — 16:9 Widescreen
# ============================================================

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# ============================================================
# PRESENTATION BUILDER
# ============================================================

class BrandedDeckBuilder:
    """Builds a branded .pptx presentation from structured JSON data."""

    def __init__(self, brand_config):
        self.brand = brand_config
        self.prs = Presentation()

        # Set 16:9 widescreen dimensions
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT

        # Brand colors
        self.primary = hex_to_rgb(brand_config["primary_color"])
        self.secondary = hex_to_rgb(brand_config["secondary_color"])
        self.accent = hex_to_rgb(brand_config["accent_color"])

        self.primary_hex = brand_config["primary_color"]
        self.secondary_hex = brand_config["secondary_color"]
        self.accent_hex = brand_config["accent_color"]

        self.light_primary = lighten_hex(brand_config["primary_color"], 0.92)
        self.light_accent = lighten_hex(brand_config["accent_color"], 0.93)
        self.dark_secondary = darken_hex(brand_config["secondary_color"], 0.2)

        # Derived colors
        self.white = RGBColor(0xFF, 0xFF, 0xFF)
        self.body_color = RGBColor(0x2A, 0x2A, 0x2A)
        self.subtle_color = RGBColor(0x77, 0x77, 0x77)
        self.dark_bg = hex_to_rgb(brand_config.get("secondary_color", "#1B1B1B"))
        self.is_dark_bg = brand_config.get("background", "dark") == "dark"

        # Resolve fonts with system-safe fallbacks
        h_font = brand_config.get("heading_font", "")
        b_font = brand_config.get("body_font", "")
        style = brand_config.get("typography_style", "modern")

        if h_font and h_font != "not specified":
            self.heading_font = FONT_FALLBACKS.get(h_font, h_font)
        else:
            self.heading_font = STYLE_FONT_MAP.get(style, ("Arial", "Calibri"))[0]

        if b_font and b_font != "not specified":
            self.body_font = FONT_FALLBACKS.get(b_font, b_font)
        else:
            self.body_font = STYLE_FONT_MAP.get(style, ("Arial", "Calibri"))[1]

        # Logo path resolution
        logo_path = brand_config.get("logo_path")
        self.logo_path = logo_path if logo_path and os.path.exists(logo_path) else None

        # Client and agency names
        self.client_name = brand_config.get("client_name", "")
        self.agency_name = brand_config.get("agency_name", "ScaleFlow")

    # ----------------------------------------------------------
    # LOW-LEVEL HELPERS
    # ----------------------------------------------------------

    def _add_blank_slide(self):
        """Add a blank slide using the blank layout (index 6)."""
        layout = self.prs.slide_layouts[6]
        return self.prs.slides.add_slide(layout)

    def _set_slide_bg(self, slide, color):
        """Set the background color of a slide."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_textbox(self, slide, left, top, width, height):
        """Add a text box to a slide and return the text frame."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        return tf

    def _set_text(self, text_frame, text, font_name=None, font_size=None,
                  color=None, bold=False, italic=False, alignment=None):
        """Set text on the first paragraph of a text frame."""
        p = text_frame.paragraphs[0]
        if alignment:
            p.alignment = alignment
        run = p.add_run()
        run.text = text
        font = run.font
        font.name = font_name or self.body_font
        if font_size:
            font.size = font_size
        if color:
            font.color.rgb = color
        font.bold = bold
        font.italic = italic
        return run

    def _add_paragraph(self, text_frame, text, font_name=None, font_size=None,
                       color=None, bold=False, italic=False, alignment=None,
                       space_before=None, space_after=None, level=0):
        """Add a new paragraph to a text frame."""
        p = text_frame.add_paragraph()
        p.level = level
        if alignment:
            p.alignment = alignment
        if space_before is not None:
            p.space_before = space_before
        if space_after is not None:
            p.space_after = space_after
        run = p.add_run()
        run.text = text
        font = run.font
        font.name = font_name or self.body_font
        if font_size:
            font.size = font_size
        if color:
            font.color.rgb = color
        font.bold = bold
        font.italic = italic
        return p

    def _add_accent_bar(self, slide, top, height=Inches(0.06)):
        """Add a horizontal accent color bar across the slide."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), top, SLIDE_WIDTH, height,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.accent
        shape.line.fill.background()
        return shape

    def _add_logo(self, slide, left, top, max_width=Inches(1.8), max_height=Inches(0.9)):
        """Add the brand logo to the slide if available. Returns the shape or None."""
        if not self.logo_path:
            return None
        try:
            img = Image.open(self.logo_path)
            img_w, img_h = img.size
            aspect = img_w / img_h

            # Fit within max bounds while maintaining aspect ratio
            width = max_width
            height = int(width / aspect)
            if height > max_height:
                height = max_height
                width = int(height * aspect)

            pic = slide.shapes.add_picture(self.logo_path, left, top, width, height)
            return pic
        except Exception:
            return None

    def _add_speaker_notes(self, slide, notes_text):
        """Add speaker notes to a slide."""
        if notes_text:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = notes_text

    # ----------------------------------------------------------
    # SLIDE: TITLE
    # ----------------------------------------------------------

    def _build_title_slide(self, slide_data):
        """Big title + subtitle, brand colors, logo, dark background."""
        slide = self._add_blank_slide()

        # Dark background
        self._set_slide_bg(slide, self.dark_bg)

        # Accent bar at top
        self._add_accent_bar(slide, Inches(0), height=Inches(0.08))

        # Title text
        tf = self._add_textbox(
            slide,
            left=Inches(1.2), top=Inches(1.8),
            width=Inches(10.5), height=Inches(2.5),
        )
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT
        self._set_text(
            tf, slide_data.get("title", ""),
            font_name=self.heading_font, font_size=Pt(54),
            color=self.primary, bold=True,
        )

        # Subtitle
        if slide_data.get("subtitle"):
            self._add_paragraph(
                tf, slide_data["subtitle"],
                font_name=self.body_font, font_size=Pt(24),
                color=self.white,
                space_before=Pt(18),
            )

        # Bottom accent bar
        self._add_accent_bar(slide, Inches(7.0), height=Inches(0.06))

        # Logo in bottom-right
        self._add_logo(
            slide,
            left=Inches(10.8), top=Inches(5.8),
            max_width=Inches(1.8), max_height=Inches(0.9),
        )

        # Agency / date footer
        footer_text = f"{self.agency_name}  |  {self.client_name}"
        ft = self._add_textbox(
            slide,
            left=Inches(1.2), top=Inches(6.5),
            width=Inches(6), height=Inches(0.5),
        )
        self._set_text(
            ft, footer_text,
            font_size=Pt(11), color=self.subtle_color,
        )

        self._add_speaker_notes(slide, slide_data.get("notes", ""))
        return slide

    # ----------------------------------------------------------
    # SLIDE: SECTION HEADER
    # ----------------------------------------------------------

    def _build_section_header_slide(self, slide_data):
        """Section divider with large centered text and accent color bar."""
        slide = self._add_blank_slide()

        # Dark background
        self._set_slide_bg(slide, self.dark_bg)

        # Accent bar at left edge (vertical stripe)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.accent
        shape.line.fill.background()

        # Centered title
        tf = self._add_textbox(
            slide,
            left=Inches(1.5), top=Inches(2.5),
            width=Inches(10), height=Inches(2.5),
        )
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT
        self._set_text(
            tf, slide_data.get("title", "").upper(),
            font_name=self.heading_font, font_size=Pt(48),
            color=self.white, bold=True,
        )

        # Horizontal accent bar below title
        self._add_accent_bar(slide, Inches(5.2), height=Inches(0.05))

        self._add_speaker_notes(slide, slide_data.get("notes", ""))
        return slide

    # ----------------------------------------------------------
    # SLIDE: CONTENT (bullets)
    # ----------------------------------------------------------

    def _build_content_slide(self, slide_data):
        """Title + bullet points body."""
        slide = self._add_blank_slide()

        if self.is_dark_bg:
            self._set_slide_bg(slide, self.dark_bg)
            title_color = self.primary
            body_text_color = self.white
        else:
            self._set_slide_bg(slide, self.white)
            title_color = self.primary
            body_text_color = self.body_color

        # Accent bar at top
        self._add_accent_bar(slide, Inches(0), height=Inches(0.05))

        # Title
        tf_title = self._add_textbox(
            slide,
            left=Inches(1.0), top=Inches(0.6),
            width=Inches(11), height=Inches(1.0),
        )
        self._set_text(
            tf_title, slide_data.get("title", ""),
            font_name=self.heading_font, font_size=Pt(36),
            color=title_color, bold=True,
        )

        # Thin accent underline below title
        underline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.0), Inches(1.55), Inches(2.5), Inches(0.04),
        )
        underline.fill.solid()
        underline.fill.fore_color.rgb = self.accent
        underline.line.fill.background()

        # Body bullets
        body = slide_data.get("body", [])
        if isinstance(body, str):
            body = [body]

        tf_body = self._add_textbox(
            slide,
            left=Inches(1.0), top=Inches(2.0),
            width=Inches(11), height=Inches(4.8),
        )
        # Remove default empty paragraph text
        tf_body.paragraphs[0].text = ""

        for i, bullet in enumerate(body):
            if i == 0:
                p = tf_body.paragraphs[0]
            else:
                p = tf_body.add_paragraph()

            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(8)
            p.space_after = Pt(8)

            # Accent bullet marker
            marker_run = p.add_run()
            marker_run.text = "\u25cf  "
            marker_run.font.name = self.body_font
            marker_run.font.size = Pt(12)
            marker_run.font.color.rgb = self.accent

            # Bullet text
            text_run = p.add_run()
            text_run.text = bullet
            text_run.font.name = self.body_font
            text_run.font.size = Pt(20)
            text_run.font.color.rgb = body_text_color

        self._add_speaker_notes(slide, slide_data.get("notes", ""))
        return slide

    # ----------------------------------------------------------
    # SLIDE: IMAGE
    # ----------------------------------------------------------

    def _build_image_slide(self, slide_data):
        """Title + image, or placeholder text if image missing."""
        slide = self._add_blank_slide()
        self._set_slide_bg(slide, self.dark_bg)

        # Accent bar at top
        self._add_accent_bar(slide, Inches(0), height=Inches(0.05))

        # Title
        tf_title = self._add_textbox(
            slide,
            left=Inches(1.0), top=Inches(0.5),
            width=Inches(11), height=Inches(0.9),
        )
        self._set_text(
            tf_title, slide_data.get("title", ""),
            font_name=self.heading_font, font_size=Pt(32),
            color=self.primary, bold=True,
        )

        image_path = slide_data.get("image_path", "")
        placeholder_text = slide_data.get("image_placeholder", "")

        if image_path and os.path.exists(image_path):
            # Insert actual image, centered and fitted
            try:
                img = Image.open(image_path)
                img_w, img_h = img.size
                aspect = img_w / img_h

                max_w = Inches(11.0)
                max_h = Inches(5.5)

                width = max_w
                height = int(width / aspect)
                if height > max_h:
                    height = max_h
                    width = int(height * aspect)

                # Center horizontally
                left = int((SLIDE_WIDTH - width) / 2)
                top = Inches(1.6)

                slide.shapes.add_picture(image_path, left, top, width, height)
            except Exception:
                self._add_image_placeholder(slide, placeholder_text or "[Image could not be loaded]")
        else:
            self._add_image_placeholder(slide, placeholder_text or "[Image not found]")

        self._add_speaker_notes(slide, slide_data.get("notes", ""))
        return slide

    def _add_image_placeholder(self, slide, text):
        """Add a styled placeholder rectangle with description text."""
        # Placeholder box
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.0), Inches(1.6),
            Inches(11.0), Inches(5.2),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(darken_hex(self.secondary_hex, 0.1))
        shape.line.color.rgb = self.accent
        shape.line.width = Pt(1.5)

        # Placeholder text inside shape
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Icon-like marker
        marker_run = tf.paragraphs[0].add_run()
        marker_run.text = "\U0001f5bc\ufe0f"
        marker_run.font.size = Pt(36)

        # Description
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(12)
        run = p.add_run()
        run.text = text
        run.font.name = self.body_font
        run.font.size = Pt(14)
        run.font.color.rgb = self.subtle_color
        run.font.italic = True

    # ----------------------------------------------------------
    # SLIDE: COMPARISON / TABLE
    # ----------------------------------------------------------

    def _build_comparison_slide(self, slide_data):
        """Title + data table."""
        slide = self._add_blank_slide()

        if self.is_dark_bg:
            self._set_slide_bg(slide, self.dark_bg)
            title_color = self.primary
        else:
            self._set_slide_bg(slide, self.white)
            title_color = self.primary

        # Accent bar at top
        self._add_accent_bar(slide, Inches(0), height=Inches(0.05))

        # Title
        tf_title = self._add_textbox(
            slide,
            left=Inches(1.0), top=Inches(0.5),
            width=Inches(11), height=Inches(0.9),
        )
        self._set_text(
            tf_title, slide_data.get("title", ""),
            font_name=self.heading_font, font_size=Pt(32),
            color=title_color, bold=True,
        )

        # Build table
        table_data = slide_data.get("table", {})
        headers = table_data.get("headers", [])
        rows = table_data.get("rows", [])

        if not headers:
            self._add_speaker_notes(slide, slide_data.get("notes", ""))
            return slide

        num_cols = len(headers)
        num_rows = len(rows) + 1  # +1 for header row

        # Calculate table dimensions
        table_width = Inches(11.0)
        table_left = Inches(1.0)
        table_top = Inches(1.7)
        row_height = Inches(0.6)
        table_height = row_height * num_rows

        # Add table shape
        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            table_left, table_top,
            table_width, table_height,
        )
        table = table_shape.table

        # Style header row
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.secondary

            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = header
            run.font.name = self.body_font
            run.font.size = Pt(14)
            run.font.color.rgb = self.white
            run.font.bold = True

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Style data rows
        light_bg = hex_to_rgb(lighten_hex(self.primary_hex, 0.92))
        for row_idx, row_data in enumerate(rows):
            for col_idx, value in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = ""

                # Alternating row colors
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = light_bg
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.white

                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = str(value)
                run.font.name = self.body_font
                run.font.size = Pt(13)
                run.font.color.rgb = self.body_color

                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Accent bottom border on table
        self._add_accent_bar(
            slide,
            table_top + table_height,
            height=Inches(0.03),
        )

        self._add_speaker_notes(slide, slide_data.get("notes", ""))
        return slide

    # ----------------------------------------------------------
    # SLIDE: TWO COLUMN
    # ----------------------------------------------------------

    def _build_two_column_slide(self, slide_data):
        """Title + two content columns."""
        slide = self._add_blank_slide()

        if self.is_dark_bg:
            self._set_slide_bg(slide, self.dark_bg)
            title_color = self.primary
            body_text_color = self.white
        else:
            self._set_slide_bg(slide, self.white)
            title_color = self.primary
            body_text_color = self.body_color

        # Accent bar at top
        self._add_accent_bar(slide, Inches(0), height=Inches(0.05))

        # Title
        tf_title = self._add_textbox(
            slide,
            left=Inches(1.0), top=Inches(0.5),
            width=Inches(11), height=Inches(0.9),
        )
        self._set_text(
            tf_title, slide_data.get("title", ""),
            font_name=self.heading_font, font_size=Pt(32),
            color=title_color, bold=True,
        )

        # Accent underline
        underline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.0), Inches(1.4), Inches(2.5), Inches(0.04),
        )
        underline.fill.solid()
        underline.fill.fore_color.rgb = self.accent
        underline.line.fill.background()

        # Left column
        left_content = slide_data.get("left", [])
        if isinstance(left_content, str):
            left_content = [left_content]

        tf_left = self._add_textbox(
            slide,
            left=Inches(1.0), top=Inches(1.9),
            width=Inches(5.2), height=Inches(4.8),
        )
        tf_left.paragraphs[0].text = ""

        for i, item in enumerate(left_content):
            if i == 0:
                p = tf_left.paragraphs[0]
            else:
                p = tf_left.add_paragraph()
            p.space_before = Pt(6)
            p.space_after = Pt(6)

            marker = p.add_run()
            marker.text = "\u25cf  "
            marker.font.size = Pt(10)
            marker.font.color.rgb = self.accent

            run = p.add_run()
            run.text = item
            run.font.name = self.body_font
            run.font.size = Pt(16)
            run.font.color.rgb = body_text_color

        # Vertical divider
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(6.55), Inches(1.9), Inches(0.03), Inches(4.5),
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = self.accent
        divider.line.fill.background()

        # Right column
        right_content = slide_data.get("right", [])
        if isinstance(right_content, str):
            right_content = [right_content]

        tf_right = self._add_textbox(
            slide,
            left=Inches(7.0), top=Inches(1.9),
            width=Inches(5.2), height=Inches(4.8),
        )
        tf_right.paragraphs[0].text = ""

        for i, item in enumerate(right_content):
            if i == 0:
                p = tf_right.paragraphs[0]
            else:
                p = tf_right.add_paragraph()
            p.space_before = Pt(6)
            p.space_after = Pt(6)

            marker = p.add_run()
            marker.text = "\u25cf  "
            marker.font.size = Pt(10)
            marker.font.color.rgb = self.accent

            run = p.add_run()
            run.text = item
            run.font.name = self.body_font
            run.font.size = Pt(16)
            run.font.color.rgb = body_text_color

        self._add_speaker_notes(slide, slide_data.get("notes", ""))
        return slide

    # ----------------------------------------------------------
    # SLIDE: CLOSING
    # ----------------------------------------------------------

    def _build_closing_slide(self, slide_data):
        """Thank you slide with brand presence."""
        slide = self._add_blank_slide()

        # Dark background
        self._set_slide_bg(slide, self.dark_bg)

        # Accent bar at top
        self._add_accent_bar(slide, Inches(0), height=Inches(0.08))

        # Main "Thank You" title
        tf = self._add_textbox(
            slide,
            left=Inches(1.0), top=Inches(2.0),
            width=Inches(11), height=Inches(2.0),
        )
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        self._set_text(
            tf, slide_data.get("title", "Thank You"),
            font_name=self.heading_font, font_size=Pt(54),
            color=self.primary, bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        # Subtitle
        if slide_data.get("subtitle"):
            self._add_paragraph(
                tf, slide_data["subtitle"],
                font_name=self.body_font, font_size=Pt(20),
                color=self.subtle_color,
                alignment=PP_ALIGN.CENTER,
                space_before=Pt(18),
            )

        # Accent bar
        self._add_accent_bar(slide, Inches(4.8), height=Inches(0.05))

        # Logo centered at bottom
        if self.logo_path:
            self._add_logo(
                slide,
                left=Inches(5.5), top=Inches(5.5),
                max_width=Inches(2.2), max_height=Inches(1.0),
            )

        # Agency footer
        ft = self._add_textbox(
            slide,
            left=Inches(1.0), top=Inches(6.6),
            width=Inches(11), height=Inches(0.5),
        )
        ft.paragraphs[0].alignment = PP_ALIGN.CENTER
        self._set_text(
            ft, f"Prepared by {self.agency_name}",
            font_size=Pt(11), color=self.subtle_color,
            alignment=PP_ALIGN.CENTER,
        )

        # Bottom accent bar
        self._add_accent_bar(slide, Inches(7.2), height=Inches(0.08))

        self._add_speaker_notes(slide, slide_data.get("notes", ""))
        return slide

    # ----------------------------------------------------------
    # LAYOUT DISPATCHER
    # ----------------------------------------------------------

    LAYOUT_MAP = {
        "title": "_build_title_slide",
        "section_header": "_build_section_header_slide",
        "content": "_build_content_slide",
        "image": "_build_image_slide",
        "comparison": "_build_comparison_slide",
        "table": "_build_comparison_slide",
        "two_column": "_build_two_column_slide",
        "closing": "_build_closing_slide",
    }

    def _build_slide(self, slide_data):
        """Dispatch to the correct slide builder based on layout type."""
        layout = slide_data.get("layout", "content")
        method_name = self.LAYOUT_MAP.get(layout)

        if method_name is None:
            print(f"Warning: Unknown layout '{layout}', falling back to 'content'.")
            method_name = "_build_content_slide"

        method = getattr(self, method_name)
        return method(slide_data)

    # ----------------------------------------------------------
    # BUILD FROM JSON
    # ----------------------------------------------------------

    def build_from_data(self, presentation_data):
        """Build the full presentation from the 'presentation' section of JSON input."""
        slides = presentation_data.get("slides", [])

        if not slides:
            print("Warning: No slides found in input data.")
            return

        for slide_data in slides:
            self._build_slide(slide_data)

        print(f"Built {len(slides)} slides.")

    # ----------------------------------------------------------
    # SAVE
    # ----------------------------------------------------------

    def save(self, output_path):
        """Save the presentation to disk."""
        out_dir = os.path.dirname(output_path)
        if out_dir:
            os.makedirs(out_dir, exist_ok=True)
        self.prs.save(output_path)
        print(f"Branded presentation saved: {output_path}")
        return output_path


# ============================================================
# CLI
# ============================================================

def main():
    parser = argparse.ArgumentParser(
        description="Generate a branded .pptx presentation from JSON data.",
        epilog="Example: python3 generate_pptx.py --input deck.json --output output.pptx",
    )
    parser.add_argument(
        "--input", required=True,
        help="Path to JSON input file containing brand config and slide data",
    )
    parser.add_argument(
        "--output", required=True,
        help="Path for the output .pptx file",
    )
    args = parser.parse_args()

    # Validate input file
    if not os.path.exists(args.input):
        print(f"Error: Input file not found: {args.input}", file=sys.stderr)
        sys.exit(1)

    # Load JSON data
    try:
        with open(args.input, "r", encoding="utf-8") as f:
            data = json.load(f)
    except json.JSONDecodeError as e:
        print(f"Error: Invalid JSON in input file: {e}", file=sys.stderr)
        sys.exit(1)

    # Validate required keys
    if "brand" not in data:
        print("Error: JSON input must contain a 'brand' object.", file=sys.stderr)
        sys.exit(1)
    if "presentation" not in data:
        print("Error: JSON input must contain a 'presentation' object.", file=sys.stderr)
        sys.exit(1)

    # Build and save
    builder = BrandedDeckBuilder(data["brand"])
    builder.build_from_data(data["presentation"])
    builder.save(args.output)


if __name__ == "__main__":
    main()
